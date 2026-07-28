from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from settings import OUTPUT_PPTX_DIR

NAVY = RGBColor(0x0F, 0x2A, 0x4A)
BLUE = RGBColor(0x1B, 0x4F, 0x91)
GRAY = RGBColor(0x6B, 0x72, 0x80)
LIGHT_BLUE = RGBColor(0xE7, 0xF0, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _add_title_slide(prs: Presentation, report_date: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(1.5))
    tf = tb.text_frame
    tf.text = "삼성화재 해외사업 뉴스 카드뉴스"
    tf.paragraphs[0].font.size = Pt(36)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = WHITE

    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.7), Inches(0.8))
    tb2.text_frame.text = f"{report_date} 생성"
    tb2.text_frame.paragraphs[0].font.size = Pt(18)
    tb2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xC7, 0xD6, 0xEA)


def _add_card_slide(prs: Presentation, country: str, article) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    y = Inches(0.35)

    if getattr(article, "is_samsung_fire_direct", False):
        badge = slide.shapes.add_shape(1, Inches(0.5), y, Inches(2.6), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(0xB9, 0x1C, 0x1C)
        badge.line.fill.background()
        bt = badge.text_frame
        bt.text = "★ 삼성화재 직접 관련"
        bt.paragraphs[0].font.size = Pt(12)
        bt.paragraphs[0].font.color.rgb = WHITE
        bt.paragraphs[0].alignment = PP_ALIGN.CENTER
        y += Inches(0.55)

    meta = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), Inches(0.4))
    meta.text_frame.text = (
        f"우선순위 {article.priority_rank} · {article.keyword_label}    |    {country}    |    "
        f"{article.local_time_display}"
    )
    meta.text_frame.paragraphs[0].font.size = Pt(12)
    meta.text_frame.paragraphs[0].font.color.rgb = GRAY
    y += Inches(0.45)

    stars = slide.shapes.add_textbox(Inches(0.5), y, Inches(4), Inches(0.4))
    stars.text_frame.text = "★" * article.importance_stars + "☆" * (5 - article.importance_stars)
    stars.text_frame.paragraphs[0].font.size = Pt(16)
    stars.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xF5, 0x9E, 0x0B)
    y += Inches(0.5)

    headline = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), Inches(1.0))
    headline.text_frame.word_wrap = True
    headline.text_frame.text = article.title
    headline.text_frame.paragraphs[0].font.size = Pt(24)
    headline.text_frame.paragraphs[0].font.bold = True
    headline.text_frame.paragraphs[0].font.color.rgb = NAVY
    y += Inches(1.1)

    summary_box = slide.shapes.add_shape(1, Inches(0.5), y, Inches(12.3), Inches(1.8))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = LIGHT_BLUE
    summary_box.line.fill.background()
    sf = summary_box.text_frame
    sf.word_wrap = True
    sf.text = article.summary
    sf.paragraphs[0].font.size = Pt(14)
    y += Inches(2.0)

    why = slide.shapes.add_textbox(Inches(0.5), y, Inches(12.3), Inches(0.8))
    why.text_frame.word_wrap = True
    why.text_frame.text = f"왜 중요한가: {article.why_important}"
    why.text_frame.paragraphs[0].font.size = Pt(13)
    y += Inches(0.9)

    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4))
    footer.text_frame.text = f"출처: {article.source_name}   |   원문: {article.url}"
    footer.text_frame.paragraphs[0].font.size = Pt(10)
    footer.text_frame.paragraphs[0].font.color.rgb = GRAY


def _add_empty_slide(prs: Presentation, country: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1))
    tb.text_frame.text = f"[{country}] 확인된 뉴스 없음"
    tb.text_frame.paragraphs[0].font.size = Pt(28)
    tb.text_frame.paragraphs[0].font.color.rgb = GRAY
    tb.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def render_pptx(articles_by_country: dict) -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    report_date = datetime.now().strftime("%Y-%m-%d")
    _add_title_slide(prs, report_date)

    for country, articles in articles_by_country.items():
        if articles:
            for a in articles:
                _add_card_slide(prs, country, a)
        else:
            _add_empty_slide(prs, country)

    out_path = OUTPUT_PPTX_DIR / f"cardnews_{report_date}.pptx"
    prs.save(str(out_path))
    return out_path
